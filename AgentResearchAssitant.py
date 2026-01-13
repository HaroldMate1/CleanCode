#!/usr/bin/env python3
"""
IAthon – Iteration 6.6 (Token Tracking Version - Fixed)
New: Comprehensive token usage logging and cost estimation.
"""

import argparse
import os
import shutil      
import subprocess
import re
import time
import warnings
from pathlib import Path
from dataclasses import dataclass, field
from typing import Dict

warnings.filterwarnings("ignore")

# --- Core scientific stack ---
import pandas as pd
import numpy as np
from scipy import stats

# --- Plotting (headless-safe) ---
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import seaborn as sns  
import plotly.express as px
import plotly.graph_objects as go

from openai import OpenAI

# Try to import python-pptx for template creation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# ======================================================
# Token Usage Tracker
# ======================================================
@dataclass
class TokenUsage:
    """Track token usage per model"""
    prompt_tokens: int = 0
    completion_tokens: int = 0
    total_tokens: int = 0
    
    def add(self, prompt: int, completion: int):
        self.prompt_tokens += prompt
        self.completion_tokens += completion
        self.total_tokens += (prompt + completion)

import json
import re
import requests
from datetime import datetime
from pathlib import Path
from typing import Dict, Optional
from dataclasses import dataclass, field
from bs4 import BeautifulSoup

# ======================================================
# FALLBACK PRICING (if web fetch fails)
# ======================================================
FALLBACK_PRICING = {
    "gpt-4o": {"input": 2.50, "output": 10.00},
    "gpt-4o-mini": {"input": 0.150, "output": 0.600},
    "o1": {"input": 15.00, "output": 60.00},
    "o1-mini": {"input": 3.00, "output": 12.00},
    "o3-mini": {"input": 1.10, "output": 4.40},
    "gpt-4-turbo": {"input": 10.00, "output": 30.00},
    "gpt-4": {"input": 30.00, "output": 60.00},
    "gpt-3.5-turbo": {"input": 0.50, "output": 1.50},
    "dall-e-3": {"standard_1024": 0.040, "standard_1792": 0.080, "hd_1024": 0.080, "hd_1792": 0.120},
    "dall-e-2": {"1024": 0.020, "512": 0.018, "256": 0.016},
}

@dataclass
class TokenUsage:
    """Track token usage per model"""
    prompt_tokens: int = 0
    completion_tokens: int = 0
    total_tokens: int = 0
    
    def add(self, prompt: int, completion: int):
        self.prompt_tokens += prompt
        self.completion_tokens += completion
        self.total_tokens += (prompt + completion)

@dataclass
class UsageTracker:
    """
    Central tracker for all API usage with live pricing fetch
    """
    models: Dict[str, TokenUsage] = field(default_factory=dict)
    pricing_cache: Dict[str, Dict[str, float]] = field(default_factory=dict)
    pricing_source: str = "unknown"
    cache_file: Path = field(default_factory=lambda: Path.home() / ".cache" / "iathlon" / "openai_pricing.json")
    
    def __post_init__(self):
        # Ensure cache directory exists
        self.cache_file.parent.mkdir(parents=True, exist_ok=True)
        
        # Try to fetch live pricing on initialization
        self.fetch_pricing()
    
    def _fetch_live_pricing_from_web(self) -> Optional[Dict]:
        """
        Attempt to fetch pricing from OpenAI's website
        Returns dict of pricing or None if failed
        """
        print("🌐 Attempting to fetch live pricing from OpenAI...")
        
        urls_to_try = [
            "https://openai.com/api/pricing/",
            "https://platform.openai.com/docs/pricing",
            "https://openai.com/pricing",
        ]
        
        for url in urls_to_try:
            try:
                response = requests.get(
                    url,
                    timeout=10,
                    headers={'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
                )
                
                if response.status_code == 200:
                    soup = BeautifulSoup(response.text, 'html.parser')
                    text = soup.get_text()
                    
                    # Extract pricing using regex patterns
                    extracted = {}
                    
                    # Pattern for model pricing: model name followed by prices
                    patterns = [
                        r'(gpt-[\w.-]+).*?\$\s*([\d.]+)\s*(?:/|per)\s*(?:1M|million).*?input.*?\$\s*([\d.]+)\s*(?:/|per)\s*(?:1M|million).*?output',
                        r'(o\d+(?:-[\w]+)?).*?\$\s*([\d.]+)\s*(?:/|per)\s*(?:1M|million).*?input.*?\$\s*([\d.]+)\s*(?:/|per)\s*(?:1M|million).*?output',
                    ]
                    
                    for pattern in patterns:
                        for match in re.finditer(pattern, text, re.IGNORECASE | re.DOTALL):
                            model = match.group(1).strip()
                            input_price = float(match.group(2))
                            output_price = float(match.group(3))
                            extracted[model] = {"input": input_price, "output": output_price}
                    
                    if extracted:
                        print(f"✅ Successfully fetched pricing for {len(extracted)} models from {url}")
                        
                        # Save to cache
                        try:
                            cache_data = {
                                'timestamp': datetime.now().isoformat(),
                                'source': url,
                                'pricing': extracted
                            }
                            with open(self.cache_file, 'w') as f:
                                json.dump(cache_data, f, indent=2)
                            print(f"💾 Saved live pricing to cache: {self.cache_file}")
                        except Exception as e:
                            print(f"⚠️  Failed to save cache: {e}")
                        
                        return extracted
            
            except Exception as e:
                continue
        
        print("⚠️  Failed to fetch live pricing from all URLs")
        return None
    
    def _load_from_cache(self) -> Optional[Dict]:
        """Load pricing from cache file"""
        try:
            if self.cache_file.exists():
                with open(self.cache_file, 'r') as f:
                    cache_data = json.load(f)
                
                timestamp = cache_data.get('timestamp', 'unknown')
                pricing = cache_data.get('pricing', {})
                
                if pricing:
                    print(f"💾 Loaded cached pricing from {timestamp}")
                    return pricing
        except Exception as e:
            print(f"⚠️  Failed to load cache: {e}")
        
        return None
    
    def fetch_pricing(self):
        """
        Fetch current pricing with fallback chain:
        1. Try live web fetch
        2. Try cached pricing
        3. Use hardcoded fallback
        """
        # Try live fetch first
        live_pricing = self._fetch_live_pricing_from_web()
        
        if live_pricing:
            self.pricing_cache = live_pricing
            self.pricing_source = "live"
            return self.pricing_cache
        
        # Try cache
        cached_pricing = self._load_from_cache()
        if cached_pricing:
            self.pricing_cache = cached_pricing
            self.pricing_source = "cache"
            return self.pricing_cache
        
        # Use fallback
        print("📋 Using hardcoded fallback pricing (January 2025)")
        self.pricing_cache = FALLBACK_PRICING.copy()
        self.pricing_source = "fallback"
        
        return self.pricing_cache
    
    def _get_model_pricing(self, model_name: str) -> Dict[str, float]:
        """Get pricing for a specific model with intelligent matching"""
        # Try exact match
        if model_name in self.pricing_cache:
            return self.pricing_cache[model_name]
        
        # Try pattern matching (e.g., gpt-4o-2024-11-20 matches gpt-4o)
        for known_model, pricing in self.pricing_cache.items():
            if model_name.startswith(known_model):
                return pricing
        
        # Infer from model name
        model_lower = model_name.lower()
        if 'gpt-5' in model_lower or 'o4' in model_lower:
            return {"input": 15.00, "output": 60.00}
        elif 'gpt-4' in model_lower or 'o1' in model_lower:
            return {"input": 30.00, "output": 60.00}
        elif 'mini' in model_lower or 'gpt-3.5' in model_lower:
            return {"input": 0.50, "output": 1.50}
        else:
            print(f"⚠️  Unknown model '{model_name}', using conservative estimate")
            return {"input": 30.00, "output": 60.00}
    
    def record(self, model: str, usage):
        """Record usage from OpenAI response"""
        if model not in self.models:
            self.models[model] = TokenUsage()
        
        prompt = 0
        completion = 0
        
        # Try different attribute names
        for prompt_attr in ['prompt_tokens', 'input_tokens', 'total_input_tokens']:
            if hasattr(usage, prompt_attr):
                prompt = getattr(usage, prompt_attr)
                break
        
        for completion_attr in ['completion_tokens', 'output_tokens', 'total_output_tokens']:
            if hasattr(usage, completion_attr):
                completion = getattr(usage, completion_attr)
                break
        
        # Fallback to dict
        if prompt == 0 and completion == 0:
            try:
                if hasattr(usage, '__dict__'):
                    usage_dict = usage.__dict__
                elif hasattr(usage, 'model_dump'):
                    usage_dict = usage.model_dump()
                else:
                    usage_dict = dict(usage)
                
                prompt = usage_dict.get('prompt_tokens', 0) or usage_dict.get('input_tokens', 0)
                completion = usage_dict.get('completion_tokens', 0) or usage_dict.get('output_tokens', 0)
            except:
                print(f"⚠️  Warning: Could not extract token usage")
        
        self.models[model].add(prompt, completion)
    
    def record_image(self, model: str, num_images: int = 1, size: str = "1024x1024", quality: str = "standard"):
        """Record image generation usage"""
        if model not in self.models:
            self.models[model] = TokenUsage()
        
        self.models[model].add(num_images, 0)
        
        if not hasattr(self.models[model], 'image_specs'):
            self.models[model].image_specs = []
        self.models[model].image_specs.append({'size': size, 'quality': quality})
    
    def print_summary(self):
        """Print detailed usage summary with cost estimates"""
        print("\n" + "="*80)
        print("💰 TOKEN USAGE & COST SUMMARY")
        print("="*80)
        
        # Show pricing source
        source_indicators = {
            'live': '✅ Live from OpenAI',
            'cache': '💾 Cached from previous fetch',
            'fallback': '📋 Hardcoded fallback (Jan 2025)'
        }
        print(f"📊 Pricing Source: {source_indicators.get(self.pricing_source, 'Unknown')}")
        print(f"📅 Report Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        print("="*80)
        
        total_cost = 0.0
        
        for model, usage in self.models.items():
            print(f"\n🤖 Model: {model}")
            
            # Check if image model
            if 'dall-e' in model.lower():
                num_images = usage.prompt_tokens
                print(f"   Images generated:  {num_images}")
                
                if hasattr(usage, 'image_specs'):
                    model_cost = 0
                    pricing = self._get_model_pricing(model)
                    
                    for spec in usage.image_specs:
                        size = spec['size']
                        quality = spec['quality']
                        
                        if 'dall-e-3' in model.lower():
                            key = f"{'hd' if quality == 'hd' else 'standard'}_{size.split('x')[0]}"
                            cost_per_image = pricing.get(key, 0.040)
                        else:
                            cost_per_image = pricing.get(size.split('x')[0], 0.020)
                        
                        model_cost += cost_per_image
                        print(f"   Size: {size}, Quality: {quality}, Cost: ${cost_per_image:.4f}")
                    
                    total_cost += model_cost
                    print(f"   Total cost:        ${model_cost:.6f}")
            else:
                # Regular token-based model
                print(f"   Input tokens:      {usage.prompt_tokens:>12,}")
                print(f"   Output tokens:     {usage.completion_tokens:>12,}")
                print(f"   Total tokens:      {usage.total_tokens:>12,}")
                
                pricing = self._get_model_pricing(model)
                input_cost = (usage.prompt_tokens / 1_000_000) * pricing["input"]
                output_cost = (usage.completion_tokens / 1_000_000) * pricing["output"]
                model_cost = input_cost + output_cost
                total_cost += model_cost
                
                print(f"   ─────────────────────────────────")
                print(f"   Input cost:        ${input_cost:>12.6f}")
                print(f"   Output cost:       ${output_cost:>12.6f}")
                print(f"   Model total:       ${model_cost:>12.6f}")
                print(f"   Rate: ${pricing['input']:.2f}/${pricing['output']:.2f} per 1M tokens")
        
        print(f"\n{'='*80}")
        print(f"💵 TOTAL ESTIMATED COST: ${total_cost:.6f}")
        print("="*80)
        
        return total_cost
    
    def save_log(self, output_dir: Path):
        """Save detailed log to file"""
        log_path = output_dir / "token_usage.log"
        
        with open(log_path, 'w') as f:
            f.write("="*80 + "\n")
            f.write("TOKEN USAGE & COST LOG\n")
            f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"Pricing Source: {self.pricing_source}\n")
            f.write("="*80 + "\n\n")
            
            total_cost = 0.0
            
            for model, usage in self.models.items():
                f.write(f"Model: {model}\n")
                
                if 'dall-e' in model.lower():
                    num_images = usage.prompt_tokens
                    f.write(f"  Images: {num_images}\n")
                    
                    if hasattr(usage, 'image_specs'):
                        model_cost = 0
                        pricing = self._get_model_pricing(model)
                        
                        for spec in usage.image_specs:
                            size = spec['size']
                            quality = spec['quality']
                            
                            if 'dall-e-3' in model.lower():
                                key = f"{'hd' if quality == 'hd' else 'standard'}_{size.split('x')[0]}"
                                cost = pricing.get(key, 0.040)
                            else:
                                cost = pricing.get(size.split('x')[0], 0.020)
                            
                            model_cost += cost
                            f.write(f"  {size} {quality}: ${cost:.4f}\n")
                        
                        total_cost += model_cost
                        f.write(f"  Total: ${model_cost:.6f}\n")
                else:
                    f.write(f"  Input tokens:  {usage.prompt_tokens:,}\n")
                    f.write(f"  Output tokens: {usage.completion_tokens:,}\n")
                    f.write(f"  Total tokens:  {usage.total_tokens:,}\n")
                    
                    pricing = self._get_model_pricing(model)
                    input_cost = (usage.prompt_tokens / 1_000_000) * pricing["input"]
                    output_cost = (usage.completion_tokens / 1_000_000) * pricing["output"]
                    model_cost = input_cost + output_cost
                    total_cost += model_cost
                    
                    f.write(f"  Input cost:  ${input_cost:.6f}\n")
                    f.write(f"  Output cost: ${output_cost:.6f}\n")
                    f.write(f"  Total cost:  ${model_cost:.6f}\n")
                
                f.write("\n")
            
            f.write("="*80 + "\n")
            f.write(f"TOTAL COST: ${total_cost:.6f}\n")
        
        print(f"📄 Token usage log saved to: {log_path}")

# ======================================================
# Utilities
# ======================================================
def get_api_key(cli_key=None):
    return cli_key or os.getenv("OPENAI_API_KEY")

def robust_rmtree(path: Path):
    """Attempt to delete a directory, handling Windows/OneDrive locks."""
    if not path.exists():
        return
    for _ in range(3):
        try:
            shutil.rmtree(path)
            return
        except PermissionError:
            time.sleep(1)
    for item in path.iterdir():
        try:
            if item.is_file(): item.unlink()
            elif item.is_dir(): shutil.rmtree(item)
        except:
            pass

def create_styled_reference(output_dir: Path, figures_dir: Path):
    """Create a styled PowerPoint reference template"""
    if not PPTX_AVAILABLE:
        print("⚠️  python-pptx not available. Install with: pip install python-pptx")
        print("   Using default PowerPoint styling instead.")
        return None
    
    print("🎨 Creating styled presentation template...")
    
    # Start with a blank presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (professional blue theme)
    DARK_BLUE = RGBColor(31, 78, 121)
    ACCENT_BLUE = RGBColor(91, 155, 213)
    GRAY = RGBColor(89, 89, 89)
    
    # Add one sample slide of each type Pandoc expects
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    
    # Customize title slide text
    title = slide1.shapes.title
    if title:
        title.text = "Sample Title"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title.text_frame.paragraphs[0].font.bold = True
    
    # Slide 2: Title and Content
    content_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_slide_layout)
    
    # Customize content slide
    title2 = slide2.shapes.title
    if title2:
        title2.text = "Sample Slide"
        title2.text_frame.paragraphs[0].font.size = Pt(32)
        title2.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        title2.text_frame.paragraphs[0].font.bold = True
    
    # Find and style the content placeholder
    for shape in slide2.placeholders:
        if shape.placeholder_format.type == 2:  # Body placeholder
            tf = shape.text_frame
            tf.text = "Sample bullet point"
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = GRAY
    
    # Save reference template
    ref_path = output_dir / "reference_template.pptx"
    prs.save(str(ref_path))
    print(f"✅ Template created: {ref_path}")
    return ref_path

# ======================================================
# 1. DATA LOADING & CLEANING
# ======================================================
def load_and_clean(path: Path):
    if path.suffix.lower() == ".csv":
        df = pd.read_csv(path)
    else:
        df = pd.read_excel(path)

    df = df.dropna(axis=1, thresh=len(df) * 0.5)
    df = df.loc[:, df.nunique() > 1]

    num_features = df.select_dtypes(include=np.number).columns.tolist()
    cat_features = df.select_dtypes(include=["object", "category"]).columns.tolist()

    return df, num_features, cat_features

# ======================================================
# 2. SAFE EXECUTION ENVIRONMENT (FIXED)
# ======================================================
# Also update CodexGenerator.generate() prompt below
class SafeRunner:
    def __init__(self, df, num_features, cat_features, output_dir: Path, verbose=False):
        self.df = df
        self.num_features = num_features
        self.cat_features = cat_features
        self.output_dir = output_dir.resolve() 
        self.verbose = verbose
        self.figures_dir = self.output_dir / "figures"
        
        if self.figures_dir.exists():
            if self.verbose: print(f"🧹 Cleaning old figures (Robust mode)...")
            robust_rmtree(self.figures_dir)
        
        self.figures_dir.mkdir(parents=True, exist_ok=True)
    
    def run(self, code):
        if "```" in code:
            match = re.search(r"```(?:python)?\n?(.*?)\n?```", code, re.DOTALL)
            if match:
                code = match.group(1)
            else:
                code = "\n".join([line for line in code.split("\n") if "```" not in line])

        # SECURITY CHECK: Block file reading operations - data is already in memory
        forbidden_patterns = [
            (r'pd\.read_csv\s*\(', "pd.read_csv() - data is already loaded in 'df' variable"),
            (r'pd\.read_excel\s*\(', "pd.read_excel() - data is already loaded in 'df' variable"),
            (r'pd\.read_json\s*\(', "pd.read_json() - data is already loaded in 'df' variable"),
            (r'pd\.read_table\s*\(', "pd.read_table() - data is already loaded in 'df' variable"),
            (r'pd\.read_sql\s*\(', "pd.read_sql() - data is already loaded in 'df' variable"),
        ]
        
        for pattern, description in forbidden_patterns:
            if re.search(pattern, code, re.IGNORECASE):
                raise RuntimeError(
                    f"❌ Code attempted to use {description}. "
                    f"The DataFrame is already available as 'df'. "
                    f"Please use the existing 'df' variable instead of reading files."
                )

        globals_safe = {
            "df": self.df, 
            "num_features": self.num_features, 
            "cat_features": self.cat_features, 
            "plt": plt, 
            "sns": sns, 
            "px": px, 
            "go": go, 
            "stats": stats, 
            "np": np, 
            "pd": pd,
            "FIGURES_DIR": str(self.figures_dir),
            "os": os  # Allow os module for path operations
        }
        
        old_cwd = os.getcwd()
        os.chdir(self.figures_dir) 
        try:
            exec(code, globals_safe)
        finally:
            os.chdir(old_cwd)


# ======================================================
# 3. DECISION MAKER (CHAT MODEL)
# ======================================================
class DecisionMaker:
    def __init__(self, api_key, tracker: UsageTracker, custom_prompt=None, verbose=False):
        self.client = OpenAI(api_key=api_key)
        self.tracker = tracker
        self.custom_prompt = custom_prompt
        self.verbose = verbose
        self.model = "gpt-5-mini"
    
    def synthesize_report(self, analysis_log, original_prompt, figures_dir, data_preview):
        if self.verbose: print("✍️ Synthesizing domain-aware report...")
        
        figs = sorted(figures_dir.glob("*.png"))
        fig_list = "\n".join([f"- {f.name}" for f in figs])

        # Build context-aware prompt
        context_section = ""
        style_section = "STYLE: High-impact scientific paper."
        
        if self.custom_prompt:
            context_section = f"""
USER CONTEXT AND REQUIREMENTS:
{self.custom_prompt}

IMPORTANT: Adapt the report to match the user's specified context, tone, length, and focus areas.
"""
            # Extract style hints if provided
            if any(word in self.custom_prompt.lower() for word in ['concise', 'brief', 'short', 'executive']):
                style_section = "STYLE: Concise executive summary format (3-5 pages max)."
            elif any(word in self.custom_prompt.lower() for word in ['detailed', 'comprehensive', 'thorough', 'in-depth']):
                style_section = "STYLE: Comprehensive technical report with detailed analysis."
            elif any(word in self.custom_prompt.lower() for word in ['casual', 'informal', 'accessible']):
                style_section = "STYLE: Accessible, conversational tone for general audience."

        synthesis_prompt = f"""
        You are a world-class Data Scientist and Technical Writer.
        
        {context_section}
        
        DATA PREVIEW (Use this to discover the context):
        {data_preview}

        ANALYSIS HISTORY:
        {analysis_log}

        AVAILABLE FIGURES:
        {fig_list}

        YOUR TASK:
        1. CONTEXT IDENTIFICATION: Based on the preview and user requirements, understand the data domain.
        2. DOMAIN VOCABULARY: Use correct terminology appropriate to the field.
        3. REPORT STRUCTURE: Intro, Detailed Results (with images), Statistical Discussion, and Conclusion.
        4. IMAGE EMBEDDING: Place ![Description](figures/filename.png) immediately after the text that discusses it.
        5. USER REQUIREMENTS: Follow any specific instructions about length, tone, focus areas, or audience.
        
        {style_section}
        """
        response = self.client.chat.completions.create(
            model=self.model, 
            messages=[{"role": "user", "content": synthesis_prompt}],
            #temperature=0.7 
        )
        
        # Track usage
        if hasattr(response, 'usage'):
            self.tracker.record(self.model, response.usage)
        
        return response.choices[0].message.content
    
    def synthesize_presentation(self, analysis_log, figures_dir, data_preview):
        """Generate presentation-optimized content with bullet points"""
        if self.verbose: print("📊 Creating presentation slides...")
        
        figs = sorted(figures_dir.glob("*.png"))
        fig_list = "\n".join([f"- {f.name}" for f in figs])

        # Adapt presentation to user context
        context_section = ""
        if self.custom_prompt:
            context_section = f"""
USER CONTEXT:
{self.custom_prompt}

Adapt the presentation style and content to match the user's requirements.
"""

        pptx_prompt = f"""
        You are creating a professional PowerPoint presentation for executives.
        
        {context_section}
        
        DATA PREVIEW:
        {data_preview}

        ANALYSIS HISTORY:
        {analysis_log}

        AVAILABLE FIGURES (CRITICAL - You MUST use these exact filenames):
        {fig_list}

        YOUR TASK - Create a slide deck in Markdown format:
        
        FIRST: Analyze the data context and identify the domain (e.g., sports, finance, healthcare, etc.)
        
        STRUCTURE (Use # for slide titles):
        
        # [Title Slide - Catchy title based on data context]
        
        [Subtitle with context]
        
        ---
        
        # Executive Summary
        
        - 3-4 key bullet points (high-level insights only)
        - Each bullet should be ONE concise line
        
        ---
        
        # Data Overview
        
        - What this data represents
        - Key metrics tracked
        - Sample size and scope
        
        ---
        
        # Key Finding 1: [Descriptive Title]
        
        - 2-3 bullet points summarizing the insight
        - Keep each bullet to ONE line
        
        ![](figures/EXACT_FIGURE_NAME_1.png)
        
        ---
        
        # Key Finding 2: [Descriptive Title]
        
        - 2-3 bullet points
        
        ![](figures/EXACT_FIGURE_NAME_2.png)
        
        ---
        
        [Continue for EACH figure in the list above - create one findings slide per figure]
        
        # Statistical Insights
        
        - 3-4 bullets highlighting statistical significance
        - Correlations or patterns discovered
        
        ---
        
        # Conclusions & Recommendations
        
        - 3-4 actionable takeaways
        - Each as a single, impactful bullet
        
        CRITICAL RULES FOR IMAGES:
        - Use EXACTLY: ![](figures/filename.png) - NO alt text in brackets
        - Use the EXACT filenames from the list above
        - Place image AFTER the bullet points on each findings slide
        - Create ONE slide per figure
        - Use --- to separate slides
        
        OTHER RULES:
        - Maximum 5 bullets per slide
        - Each bullet must be ONE line (max 15 words)
        - Use domain-appropriate terminology
        - Be specific with numbers when relevant
        - NO paragraphs, NO long explanations
        """
        
        response = self.client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": pptx_prompt}],
            #temperature=0.6
        )
        
        # Track usage
        if hasattr(response, 'usage'):
            self.tracker.record(self.model, response.usage)
        
        content = response.choices[0].message.content
        
        # Extract domain context for image generation
        domain_match = re.search(r'#\s*(.+?)(?:\n|$)', content)
        if domain_match:
            title = domain_match.group(1).strip()
            # Generate a short domain description from the title
            domain_context = f"professional, minimalist illustration representing {title}"
        else:
            domain_context = "professional data analysis presentation"
        
        return content, domain_context
    
    def generate_title_image(self, domain_context: str, output_path: Path):
        """Generate a title slide image using DALL-E"""
        if self.verbose: print(f"🎨 Generating title image: {domain_context}")
        
        # Enhance with custom context if available
        if self.custom_prompt:
            domain_context = f"{domain_context}. Context: {self.custom_prompt[:100]}"
        
        # Create a refined prompt for DALL-E
        image_prompt = f"""
        Create a professional, modern, minimalist cover image for a business presentation.
        Theme: {domain_context}
        Style: Clean, corporate, high-quality photography or digital art
        Composition: Centered focal point with negative space for text overlay
        Colors: Professional color palette with blue/gray tones
        NO text, NO numbers, NO labels in the image
        """
        
        try:
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=image_prompt,
                size="1792x1024",  # Wide format for presentation
                quality="standard",
                n=1
            )
            
            # Track image generation usage
            self.tracker.record_image("dall-e-3", num_images=1, size="1792x1024", quality="standard")
            
            # Download the image
            import urllib.request
            image_url = response.data[0].url
            urllib.request.urlretrieve(image_url, output_path)
            
            print(f"✅ Title image saved to: {output_path}")
            return output_path
            
        except Exception as e:
            print(f"⚠️  Failed to generate title image: {e}")
            return None

    def decide(self, observations: str, step_num: int, max_steps: int, 
           num_features: list, cat_features: list, data_preview: str) -> str:
        
        # Adapt analysis approach based on custom prompt
        context_section = ""
        if self.custom_prompt:
            context_section = f"""
USER PROVIDED CONTEXT:
{self.custom_prompt}

Use this information to guide your analysis focus and methodology.
"""
        
        prompt = f"""Analyze the DATA PREVIEW below to identify the domain.
                
                {context_section}
                
                DATA PREVIEW:
                {data_preview}

                STEP: {step_num + 1} of {max_steps}
                CURRENT OBSERVATIONS: {observations}

                MISSION: 
                - Deduce the origin/subject of the data (consider user context if provided).
                - Propose a specific visualization or statistical test.
                - If columns represent times (Split/Swim/Bike), suggest converting them to seconds.
                - Align analysis with user requirements and priorities.
                
                DECISION (ONE ACTION):"""
        
        response = self.client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            #temperature=0.3,
        )
        
        # Track usage
        if hasattr(response, 'usage'):
            self.tracker.record(self.model, response.usage)
        
        return response.choices[0].message.content.strip()

# ======================================================
# 4. CODE GENERATOR (CODEX)
# ======================================================
class CodexGenerator:
    def __init__(self, api_key, tracker: UsageTracker, verbose=False):
        self.client = OpenAI(api_key=api_key)
        self.tracker = tracker
        self.verbose = verbose
        self.model = "gpt-5.1-codex-mini"

    def generate(self, instruction, num_features, cat_features, data_preview) -> str:
        prompt = f"""
            You are an expert data scientist.
            DATA PREVIEW: {data_preview}
            TASK: {instruction}
            
            CRITICAL RULES FOR DATA TYPE HANDLING:
            - ALWAYS check column dtypes before plotting: df[col].dtype
            - If a column contains time strings (e.g., HH:MM:SS), use pd.to_timedelta() and .dt.total_seconds()
            - If a column is 'object' dtype but should be numeric, use pd.to_numeric(df[col], errors='coerce')
            - Before ANY plotting operation, ensure numeric columns are actually numeric: df[col] = pd.to_numeric(df[col], errors='coerce')
            - Drop NaN values after conversion: df = df.dropna(subset=[col])
            - For categorical data in plots, convert to string explicitly: df[col].astype(str)
            
            EXAMPLE SAFE PATTERN:
            ```python
            # Ensure numeric
            df['column'] = pd.to_numeric(df['column'], errors='coerce')
            df = df.dropna(subset=['column'])
            
            # Now plot
            plt.figure(figsize=(10,6))
            plt.plot(df['column'])
            plt.savefig(os.path.join(FIGURES_DIR, 'plot.png'))
            plt.close()
            ```
            
            OTHER RULES:
            - OUTPUT PYTHON CODE ONLY. NO Markdown.
            - ALWAYS save to FIGURES_DIR.
            - Use descriptive filenames for figures.
            """
        response = self.client.responses.create(
            model=self.model,
            input=prompt,
        )
        
        # Track usage - try to access usage attribute safely
        if hasattr(response, 'usage') and response.usage is not None:
            self.tracker.record(self.model, response.usage)
        
        return response.output_text.strip()

# ======================================================
# 5. REACT ORCHESTRATOR
# ======================================================
class ReActAnalyzer:
    def __init__(self, runner, decider, coder, tracker: UsageTracker, max_steps=20, verbose=False):
        self.runner = runner
        self.decider = decider
        self.coder = coder
        self.tracker = tracker
        self.max_steps = max_steps
        self.verbose = verbose
        self.observations = []
        self.analysis_log = []
        self.error_history = []  # Track errors to avoid repeating them
        self.data_preview = self.runner.df.head(10).to_string()

    def observe(self):
        figs = sorted(self.runner.figures_dir.glob("*.png"))
        obs = f"- Figures: {len(figs)}\n- Files: " + ", ".join([f.name for f in figs]) if figs else "- No figs."
        self.observations.append(obs)

    def run(self, user_requirements): 
        print(f"\n🔬 STARTING ANALYSIS (Iteration 6.6 - Token Tracking)\n{'='*70}")
        self.observe()
        
        for step in range(self.max_steps):
            print(f"\n🔄 STEP {step + 1}/{self.max_steps}")
            
            # Include error history in context
            error_context = ""
            if self.error_history:
                recent_errors = self.error_history[-3:]  # Last 3 errors
                error_context = "\n\nRECENT ERRORS TO AVOID:\n" + "\n".join([
                    f"- Error: {e['error']}\n  From action: {e['action'][:100]}..."
                    for e in recent_errors
                ])
            
            decision = self.decider.decide(
                "\n".join(self.observations) + error_context, 
                step, 
                self.max_steps, 
                self.runner.num_features, 
                self.runner.cat_features, 
                self.data_preview
            )
            
            if "STOP" in decision.upper(): 
                print("🛑 Analysis complete (STOP signal received)")
                break
            
            code = self.coder.generate(decision, self.runner.num_features, self.runner.cat_features, self.data_preview)
            
            try:
                self.runner.run(code)
                if self.verbose: print("⚡ Success.")
                # Clear error on success
                self.observe()
                self.analysis_log.append({
                    "step": step + 1, 
                    "thought": decision, 
                    "result": self.observations[-1],
                    "status": "success"
                })
            except Exception as e:
                error_msg = str(e)
                print(f"❌ Error: {error_msg}")
                
                # Store error for context
                self.error_history.append({
                    "step": step + 1,
                    "action": decision,
                    "error": error_msg
                })
                
                # Add error to observations
                error_obs = f"- ERROR in step {step + 1}: {error_msg}"
                self.observations.append(error_obs)
                
                self.analysis_log.append({
                    "step": step + 1, 
                    "thought": decision, 
                    "result": error_obs,
                    "status": "error"
                })

        report_text = self.decider.synthesize_report(self.analysis_log, user_requirements, self.runner.figures_dir, self.data_preview)
        report_path = self.runner.output_dir / "report.md" 
        report_path.write_text(report_text, encoding='utf-8')
        print(f"✨ Analysis complete. Report saved to {report_path}")
        
        return report_path, report_text

# ======================================================
# 6. CLI
# ======================================================
def main():
    parser = argparse.ArgumentParser(
        description="IAthon - Intelligent Automated Data Analysis with customizable reports"
    )
    parser.add_argument("-i", "--input", required=True, help="Input data file (CSV or Excel)")
    parser.add_argument("-o", "--output", required=True, help="Output file path")
    parser.add_argument("-f", "--format", choices=["pdf", "docx", "pptx"], help="Output format")
    parser.add_argument("--api-key", required=False, help="OpenAI API key (or set OPENAI_API_KEY env var)")
    parser.add_argument("-v", "--verbose", action="store_true", help="Verbose output")
    parser.add_argument(
        "-p", "--prompt", 
        required=False, 
        help="Custom prompt describing data context and report requirements (e.g., 'This is triathlon race data. Create a concise executive summary with focus on performance trends.')"
    )
    parser.add_argument(
        "--prompt-file",
        required=False,
        help="Path to text file containing custom prompt"
    )
    args = parser.parse_args()

    api_key = get_api_key(args.api_key)
    input_path = Path(args.input)
    output_path = Path(args.output)
    output_dir = output_path.parent

    # Load custom prompt if provided
    custom_prompt = None
    if args.prompt_file:
        prompt_file = Path(args.prompt_file)
        if prompt_file.exists():
            custom_prompt = prompt_file.read_text(encoding='utf-8')
            print(f"📝 Loaded custom prompt from: {prompt_file}")
        else:
            print(f"⚠️  Prompt file not found: {prompt_file}")
    elif args.prompt:
        custom_prompt = args.prompt
        print(f"📝 Using custom prompt: {custom_prompt[:100]}...")

    # Initialize usage tracker
    tracker = UsageTracker()

    df, num_features, cat_features = load_and_clean(input_path)
    runner = SafeRunner(df, num_features, cat_features, output_dir, verbose=args.verbose)
    decider = DecisionMaker(api_key, tracker, custom_prompt=custom_prompt, verbose=args.verbose)
    coder = CodexGenerator(api_key, tracker, verbose=args.verbose)

    analyzer = ReActAnalyzer(runner, decider, coder, tracker, max_steps=5, verbose=args.verbose)
    
    # Use custom prompt if provided, otherwise use default
    if custom_prompt:
        user_requirements = custom_prompt
    else:
        user_requirements = "Comprehensive domain-specific report with embedded figures."
    
    report_md_path, report_text = analyzer.run(user_requirements) 

    if args.format:
        print(f"📦 Converting to {args.format}...")
        out_file = output_path.with_suffix(f".{args.format}")
        
        if args.format == "pptx":
            # Generate presentation-optimized content
            pptx_content, domain_context = decider.synthesize_presentation(
                analyzer.analysis_log, 
                runner.figures_dir, 
                analyzer.data_preview
            )
            
            # Generate title image with DALL-E
            title_image_path = runner.figures_dir / "title_image.png"
            generated_image = decider.generate_title_image(domain_context, title_image_path)
            
            # If image was generated, add it to the title slide
            if generated_image:
                # Insert image reference after the first slide
                lines = pptx_content.split('\n')
                insert_index = None
                for i, line in enumerate(lines):
                    if line.strip() == '---' and insert_index is None:
                        insert_index = i
                        break
                
                if insert_index:
                    lines.insert(insert_index, '\n![](figures/title_image.png)\n')
                    pptx_content = '\n'.join(lines)
            
            pptx_md_path = output_dir / "presentation.md"
            pptx_md_path.write_text(pptx_content, encoding='utf-8')
            print(f"📊 Presentation markdown saved to {pptx_md_path}")
            
            # Create reference document with styling
            ref_path = create_styled_reference(output_dir, runner.figures_dir)
            
            # Build Pandoc command
            pandoc_cmd = ["pandoc"]
            
            # Input file (must be relative to cwd)
            pandoc_cmd.extend([str(pptx_md_path.name)])
            
            # Output format and file
            pandoc_cmd.extend(["-t", "pptx", "-o", str(out_file.name)])
            
            # Add reference doc if template was created successfully
            if ref_path and ref_path.exists():
                pandoc_cmd.extend([f"--reference-doc={ref_path.name}"])
            
            # Resource path for images (absolute path is safer)
            pandoc_cmd.extend([f"--resource-path=.:figures:{runner.figures_dir.absolute()}"])
            
            # Add standalone flag to ensure complete document
            pandoc_cmd.append("--standalone")
            
            print(f"🔧 Running: {' '.join(pandoc_cmd)}")
            
            try:
                result = subprocess.run(
                    pandoc_cmd, 
                    cwd=output_dir, 
                    check=True,
                    capture_output=True,
                    text=True
                )
                print(f"✅ Presentation saved to {out_file}")
            except subprocess.CalledProcessError as e:
                print(f"❌ Pandoc error: {e}")
                print(f"   stdout: {e.stdout}")
                print(f"   stderr: {e.stderr}")
                print(f"   Trying without reference doc...")
                
                # Fallback: try without reference doc
                pandoc_cmd_simple = [
                    "pandoc",
                    str(pptx_md_path.name),
                    "-t", "pptx",
                    "-o", str(out_file.name),
                    f"--resource-path=.:figures:{runner.figures_dir.absolute()}",
                    "--standalone"
                ]
                subprocess.run(pandoc_cmd_simple, cwd=output_dir, check=True)
                print(f"✅ Presentation saved to {out_file} (with default styling)")
        else:
            # For PDF and DOCX, use the full report
            subprocess.run(
                ["pandoc", report_md_path.name, "-o", out_file.name],
                cwd=output_dir
            )
            print(f"✅ Document saved to {out_file}")
    
    # Print token usage LAST (after all processing is complete)
    print("\n" + "="*70)
    tracker.print_summary()
    tracker.save_log(output_dir)

if __name__ == "__main__":
    main()